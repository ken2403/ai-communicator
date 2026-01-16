#!/usr/bin/env python3
"""
Generate PowerPoint slides from markdown design files.

This script reads a markdown design file and generates a PowerPoint presentation
using the specified template. Colors are extracted from the template's theme.

Usage:
    uv run python slides/scripts/generate_pptx.py <design_file> [--template <template_name>]

Example:
    uv run python slides/scripts/generate_pptx.py slides/design/2026-01-16_warehouse-system-proposal.md
"""

import argparse
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


class TemplateColors:
    """Extract and manage colors from PowerPoint template theme."""

    def __init__(self, template_path: str):
        self.colors = self._extract_theme_colors(template_path)
        self._setup_semantic_colors()

    def _extract_theme_colors(self, template_path: str) -> dict:
        """Extract colors from template's theme XML."""
        colors = {}
        with zipfile.ZipFile(template_path, 'r') as z:
            for tf in z.namelist():
                if 'theme1.xml' in tf:
                    content = z.read(tf)
                    root = ET.fromstring(content)
                    for elem in root.iter():
                        if 'clrScheme' in elem.tag:
                            for child in elem:
                                tag = child.tag.split('}')[-1]
                                for color_elem in child:
                                    color_tag = color_elem.tag.split('}')[-1]
                                    if color_tag == 'srgbClr':
                                        val = color_elem.get('val')
                                        colors[tag] = RGBColor.from_string(val)
                                    elif color_tag == 'sysClr':
                                        val = color_elem.get('lastClr')
                                        if val:
                                            colors[tag] = RGBColor.from_string(val)
        return colors

    def _setup_semantic_colors(self):
        """Map theme colors to semantic names."""
        self.dark_navy = self.colors.get('dk1', RGBColor(0, 23, 34))
        self.gray = self.colors.get('lt1', RGBColor(151, 154, 155))
        self.light_gray = self.colors.get('dk2', RGBColor(243, 244, 244))
        self.gold = self.colors.get('lt2', RGBColor(192, 161, 113))
        self.dark_gray_blue = self.colors.get('accent1', RGBColor(130, 145, 155))
        self.mauve = self.colors.get('accent2', RGBColor(157, 120, 140))
        self.sage_green = self.colors.get('accent3', RGBColor(163, 171, 145))
        self.beige = self.colors.get('accent4', RGBColor(225, 210, 187))
        self.light_gray_blue = self.colors.get('accent5', RGBColor(209, 214, 217))
        self.burgundy = self.colors.get('accent6', RGBColor(223, 51, 72))
        self.white = RGBColor(255, 255, 255)


class SlideGenerator:
    """Generate PowerPoint slides with proper positioning."""

    # Slide dimensions for GENDA template (20" x 11.25")
    SLIDE_WIDTH = 20.0
    SLIDE_HEIGHT = 11.25

    # Content area margins (avoid right side copyright area)
    MARGIN_LEFT = 1.0
    MARGIN_RIGHT = 2.0  # Increased to avoid © GENDA Inc.
    MARGIN_TOP = 2.8
    MARGIN_BOTTOM = 1.0

    def __init__(self, template_path: str):
        self.template_path = template_path
        self.colors = TemplateColors(template_path)
        self.prs = None

        # Calculate content area
        self.content_width = self.SLIDE_WIDTH - self.MARGIN_LEFT - self.MARGIN_RIGHT
        self.content_height = self.SLIDE_HEIGHT - self.MARGIN_TOP - self.MARGIN_BOTTOM
        self.content_left = self.MARGIN_LEFT
        self.content_right = self.SLIDE_WIDTH - self.MARGIN_RIGHT

    def load_template(self):
        """Load the PowerPoint template."""
        self.prs = Presentation(self.template_path)

    def delete_all_slides(self):
        """Delete all existing slides from the presentation."""
        slide_ids = list(self.prs.slides._sldIdLst)
        for slide_id in slide_ids:
            rId = slide_id.rId
            self.prs.part.drop_rel(rId)
            self.prs.slides._sldIdLst.remove(slide_id)

    def center_left(self, total_width: float) -> float:
        """Calculate left position to center content horizontally."""
        return self.content_left + (self.content_width - total_width) / 2

    def add_title_slide(self, title: str, subtitle: str, date: str = "2026.01.XX"):
        """Add a title slide using Layout 0."""
        layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(layout)

        for shape in slide.shapes:
            if shape.has_text_frame:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type is not None:
                    ph_type = shape.placeholder_format.type.real
                    if ph_type == 1:  # TITLE
                        shape.text_frame.paragraphs[0].text = title
                        for para in shape.text_frame.paragraphs:
                            para.font.color.rgb = self.colors.gold
                            para.font.size = Pt(60)
                            para.font.bold = True
                    elif ph_type == 4:  # SUBTITLE
                        if shape.top.inches > 2.5:
                            shape.text_frame.paragraphs[0].text = subtitle
                            for para in shape.text_frame.paragraphs:
                                para.font.color.rgb = self.colors.dark_gray_blue
                                para.font.size = Pt(24)
                        else:
                            shape.text_frame.paragraphs[0].text = ""

        # Add date
        txBox = slide.shapes.add_textbox(Inches(0.83), Inches(4.0), Inches(16), Inches(0.5))
        p = txBox.text_frame.paragraphs[0]
        p.text = date
        p.font.size = Pt(20)
        p.font.color.rgb = self.colors.dark_gray_blue

        return slide

    def add_content_slide(self, title: str):
        """Add a content slide using Layout 2 with proper title."""
        layout = self.prs.slide_layouts[2]
        slide = self.prs.slides.add_slide(layout)

        for shape in slide.shapes:
            if shape.has_text_frame and hasattr(shape, 'placeholder_format'):
                try:
                    ph_type = shape.placeholder_format.type.real
                    if ph_type == 1:  # TITLE
                        shape.text_frame.paragraphs[0].text = title
                        for para in shape.text_frame.paragraphs:
                            para.font.color.rgb = self.colors.gold
                            para.font.size = Pt(36)
                            para.font.bold = True
                    elif ph_type == 4:  # SUBTITLE - clear it
                        shape.text_frame.paragraphs[0].text = ""
                except:
                    pass

        return slide

    def add_rounded_box(self, slide, left: float, top: float, width: float, height: float,
                        fill_color, text: str = "", font_size: int = 16, font_color=None):
        """Add a rounded rectangle with text."""
        if font_color is None:
            font_color = self.colors.white

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()

        if hasattr(shape, 'adjustments') and len(shape.adjustments) > 0:
            shape.adjustments[0] = 0.1

        if text:
            tf = shape.text_frame
            tf.word_wrap = True
            tf.anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.color.rgb = font_color
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

        return shape

    def add_text_box(self, slide, left: float, top: float, width: float, height: float,
                     text: str, font_size: int = 14, font_color=None, bold: bool = False,
                     align=PP_ALIGN.LEFT):
        """Add a text box."""
        if font_color is None:
            font_color = self.colors.dark_navy

        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = align
        return txBox

    def add_multiline_box(self, slide, left: float, top: float, width: float, height: float,
                          fill_color, title: str, subtitle: str = "", font_color=None,
                          title_size: int = 28, subtitle_size: int = 18):
        """Add a rounded box with title and subtitle."""
        if font_color is None:
            font_color = self.colors.white

        shape = self.add_rounded_box(slide, left, top, width, height, fill_color, "", 18, font_color)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(title_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(subtitle_size)
            p2.font.color.rgb = font_color
            p2.font.bold = False
            p2.alignment = PP_ALIGN.CENTER

        return shape

    def save(self, output_path: str):
        """Save the presentation."""
        self.prs.save(output_path)


def main():
    parser = argparse.ArgumentParser(description='Generate PowerPoint from markdown design')
    parser.add_argument('design_file', help='Path to the markdown design file')
    parser.add_argument('--template', '-t', default='genda', help='Template name (default: genda)')
    parser.add_argument('--delete-existing', action='store_true', default=True,
                        help='Delete existing slides from template (default: True)')
    args = parser.parse_args()

    # Determine paths
    design_path = Path(args.design_file)
    template_path = Path(f'./slides/templates/{args.template}.pptx')
    output_path = Path('./slides/output') / design_path.with_suffix('.pptx').name

    print(f"Design file: {design_path}")
    print(f"Template: {template_path}")
    print(f"Output: {output_path}")

    # Generate slides
    gen = SlideGenerator(str(template_path))
    gen.load_template()

    if args.delete_existing:
        gen.delete_all_slides()
        print("Deleted existing slides from template")

    # Print loaded colors
    print(f"\nLoaded theme colors:")
    for name, color in gen.colors.colors.items():
        print(f"  {name}: #{color}")

    print(f"\nContent area: {gen.content_left}in - {gen.content_right}in (width: {gen.content_width}in)")

    # Note: Actual slide generation would parse the markdown file
    # For now, this is a base class that can be extended
    print("\nBase generator ready. Extend this class for specific presentations.")


if __name__ == "__main__":
    main()
